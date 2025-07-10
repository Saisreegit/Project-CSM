import os
from werkzeug.utils import secure_filename
from datetime import datetime
import openpyxl
from docx import Document
from pptx import Presentation
import smtplib
from email.mime.text import MIMEText
import difflib
from flask import flash # Import flash for use in helper functions

# These are global variables in app.py, so we need to pass them or define them if functions
# are truly independent. For now, we'll assume audit_log is passed.
# In a larger app, you might consider a logging system or a database for audit logs.

def log_action(audit_log_list, entry):
    """Adds an entry to the in-memory audit log."""
    audit_log_list.append({'timestamp': datetime.now().isoformat(), 'entry': entry})
    print(f"LOG: {entry}") # Also print to console for immediate feedback

def detect_type(filename):
    """Detects file type based on extension."""
    ext = filename.split('.')[-1].lower()
    return {'xlsx': 'excel', 'docx': 'word', 'pptx': 'ppt', 'pdf': 'pdf'}.get(ext, 'unknown')

def save_file(file, upload_folder):
    """Saves an uploaded file to the UPLOAD_FOLDER."""
    fname = secure_filename(file.filename)
    path = os.path.join(upload_folder, fname)
    file.save(path)
    return fname, path

def send_email(to_email, subject, body, smtp_server, smtp_port, email_sender, email_password):
    """Sends an email using the configured SMTP server."""
    msg = MIMEText(body)
    msg['Subject'] = subject
    msg['From'] = email_sender
    msg['To'] = to_email
    try:
        with smtplib.SMTP(smtp_server, smtp_port) as server:
            server.starttls() # Enable TLS encryption
            server.login(email_sender, email_password)
            server.send_message(msg)
        print(f"Email sent successfully to {to_email}")
        # log_action is called within app.py, so we don't call it here to avoid circular dependency
    except smtplib.SMTPAuthenticationError:
        print("Email sending failed: Authentication Error. Check your Gmail App Password.")
        flash("Email sending failed: Authentication error. Check your email configuration (especially Gmail App Password).", "error")
    except Exception as e:
        print(f"Email sending failed: {e}")
        flash(f"Email sending failed: {e}", "error")

# --- Document Reading Functions (for preview) ---
def read_excel_sheets_with_names(path):
    """Reads all sheets from an Excel file and returns data as a dict of lists of lists,
       along with a list of sheet names."""
    try:
        wb = openpyxl.load_workbook(path)
        data = {}
        for name in wb.sheetnames:
            sheet_data = []
            for row in wb[name].iter_rows():
                row_values = [cell.value for cell in row]
                sheet_data.append(row_values)
            if sheet_data:
                data[name] = sheet_data
        return data, wb.sheetnames
    except Exception as e:
        print(f"Error reading Excel file {path}: {e}")
        flash(f"Error reading Excel file: {e}", "error")
        return {}, []

def read_docx(path):
    """Reads text from a Word document as a list of paragraph strings."""
    try:
        doc = Document(path)
        return [p.text for p in doc.paragraphs] # Keep empty paragraphs for better diff alignment
    except Exception as e:
        print(f"Error reading Word file {path}: {e}")
        flash(f"Error reading Word file: {e}", "error")
        return []

def read_ppt(path):
    """Reads text from a PowerPoint presentation."""
    try:
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    if shape.text_frame.text.strip():
                        slide_text.append(shape.text_frame.text)
            if slide_text: # Only add slide if it has text
                slides.append({'slide_number': i + 1, 'text': slide_text})
        return slides
    except Exception as e:
        print(f"Error reading PowerPoint file {path}: {e}")
        flash(f"Error reading PowerPoint file: {e}", "error")
        return []

# --- Side-by-Side Comparison Logic ---

def _get_excel_diff_html(old_excel_data, new_excel_data, selected_sheet_a=None, selected_sheet_b=None):
    """
    Compares two Excel data structures (dict of sheets) and returns
    HTML for side-by-side display with highlighting.
    Highlights changed cells in green (new) and red (deleted/modified from old).
    Supports sheet selection.
    """
    left_html = ""
    right_html = ""

    # Determine which sheets to display/compare
    sheet_a = old_excel_data.get(selected_sheet_a, []) if selected_sheet_a else (list(old_excel_data.values())[0] if old_excel_data else [])
    sheet_b = new_excel_data.get(selected_sheet_b, []) if selected_sheet_b else (list(new_excel_data.values())[0] if new_excel_data else [])

    sheet_name_a = selected_sheet_a if selected_sheet_a else (list(old_excel_data.keys())[0] if old_excel_data else "N/A")
    sheet_name_b = selected_sheet_b if selected_sheet_b else (list(new_excel_data.keys())[0] if new_excel_data else "N/A")

    left_html += f"<h3>Sheet: {sheet_name_a}</h3><div class='excel-sheet-wrapper'><table class='excel-table'>"
    right_html += f"<h3>Sheet: {sheet_name_b}</h3><div class='excel-sheet-wrapper'><table class='excel-table'>"

    max_rows = max(len(sheet_a), len(sheet_b))
    max_cols = 0
    if sheet_a: max_cols = max(max_cols, max(len(row) for row in sheet_a))
    if sheet_b: max_cols = max(max_cols, max(len(row) for row in sheet_b))

    for r_idx in range(max_rows):
        left_html += "<tr>"
        right_html += "<tr>"
        for c_idx in range(max_cols):
            old_cell_value = sheet_a[r_idx][c_idx] if r_idx < len(sheet_a) and c_idx < len(sheet_a[r_idx]) else None
            new_cell_value = sheet_b[r_idx][c_idx] if r_idx < len(sheet_b) and c_idx < len(sheet_b[r_idx]) else None

            old_display = "" if old_cell_value is None else str(old_cell_value)
            new_display = "" if new_cell_value is None else str(new_cell_value)

            cell_style_left = ""
            cell_style_right = ""

            # Highlight logic for the selected sheets
            if old_cell_value != new_cell_value:
                if old_cell_value is not None and new_cell_value is None: # Deleted from old
                    cell_style_left = 'background-color: #ffcccc;' # Red
                elif old_cell_value is None and new_cell_value is not None: # Added in new
                    cell_style_right = 'background-color: #ccffcc;' # Green
                elif old_cell_value is not None and new_cell_value is not None: # Modified
                    cell_style_left = 'background-color: #ffcccc;' # Red in original for modification
                    cell_style_right = 'background-color: #ccffcc;' # Green in new for modification

            left_html += f"<td style='{cell_style_left}'>{old_display}</td>"
            right_html += f"<td style='{cell_style_right}'>{new_display}</td>"

        left_html += "</tr>"
        right_html += "</tr>"
    left_html += "</table></div>"
    right_html += "</table></div>"

    return left_html, right_html


def _get_word_diff_html(old_paragraphs, new_paragraphs):
    """
    Compares two lists of paragraphs and returns HTML for side-by-side display
    with line-based highlighting.
    Red for deleted lines in left pane, Green for added lines in right pane.
    Modified lines will appear as a red deleted line on left and green added line on right.
    """
    differ = difflib.Differ()
    diff = list(differ.compare(old_paragraphs, new_paragraphs))

    left_html_lines = []
    right_html_lines = []

    for line in diff:
        if line.startswith('  '): # Unchanged
            left_html_lines.append(f'<span>{line[2:]}</span>')
            right_html_lines.append(f'<span>{line[2:]}</span>')
        elif line.startswith('- '): # Deleted from original
            left_html_lines.append(f'<span style="background-color: #ffcccc;">{line[2:]}</span>')
            right_html_lines.append(f'<span style="visibility: hidden;">{line[2:]}</span>') # Placeholder for alignment
        elif line.startswith('+ '): # Added to new
            left_html_lines.append(f'<span style="visibility: hidden;">{line[2:]}</span>') # Placeholder for alignment
            right_html_lines.append(f'<span style="background-color: #ccffcc;">{line[2:]}</span>')
        # '?' lines are internal to difflib, ignore for display

    return "<div class='word-preview-content'>" + "<br>".join(left_html_lines) + "</div>", \
           "<div class='word-preview-content'>" + "<br>".join(right_html_lines) + "</div>"


def compare_documents_side_by_side(doc_type, path_a, path_b, selected_sheet_a=None, selected_sheet_b=None):
    """
    Generates HTML for side-by-side comparison.
    Returns (left_panel_html, right_panel_html).
    """
    if doc_type == 'excel':
        data_a, _ = read_excel_sheets_with_names(path_a)
        data_b, _ = read_excel_sheets_with_names(path_b)
        return _get_excel_diff_html(data_a, data_b, selected_sheet_a, selected_sheet_b)
    elif doc_type == 'word':
        data_a = read_docx(path_a)
        data_b = read_docx(path_b)
        return _get_word_diff_html(data_a, data_b)
    elif doc_type == 'pdf':
        # For PDF, we can only display embeds, no direct content diff
        # We need to pass url_for or get it from app context, for now assuming it's available or path is direct
        # This part assumes Flask's url_for is accessible or you're providing full URLs
        return f"<embed src='/uploads/{os.path.basename(path_a)}' type='application/pdf' width='100%' height='600px' />", \
               f"<embed src='/uploads/{os.path.basename(path_b)}' type='application/pdf' width='100%' height='600px' />"
    elif doc_type == 'ppt':
        # PPT is too complex for simple diff, just show original previews
        data_a = read_ppt(path_a)
        data_b = read_ppt(path_b)
        left_html = "<h3>PowerPoint Preview</h3>"
        for slide in data_a:
            left_html += f"<div class='ppt-slide'><h4>Slide {slide['slide_number']}:</h4>"
            for text_line in slide['text']: left_html += f"<p>{text_line}</p>"
            left_html += "</div>"
        right_html = "<h3>PowerPoint Preview</h3>"
        for slide in data_b:
            right_html += f"<div class='ppt-slide'><h4>Slide {slide['slide_number']}:</h4>"
            for text_line in slide['text']: right_html += f"<p>{text_line}</p>"
            right_html += "</div>"
        return left_html, right_html
    else:
        return "<p>No preview or comparison available.</p>", "<p>No preview or comparison available.</p>"